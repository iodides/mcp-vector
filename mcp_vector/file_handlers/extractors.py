"""
File content extractors for different file types
"""
import os
import logging
from pathlib import Path
from typing import Optional, Dict, List, Any, Tuple

# File type handling libraries
try:
    import pypdf
    PYPDF_AVAILABLE = True
except ImportError:
    PYPDF_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import openpyxl
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

logger = logging.getLogger(__name__)

class FileHandler:
    """Base class for file handlers"""
    
    @staticmethod
    def can_handle(file_path: str) -> bool:
        """Check if this handler can process the given file"""
        raise NotImplementedError
    
    @staticmethod
    def extract_text(file_path: str) -> Tuple[str, Dict[str, Any]]:
        """
        Extract text content from file
        
        Args:
            file_path: Path to the file
            
        Returns:
            Tuple of (extracted text, metadata)
        """
        raise NotImplementedError
    
    @staticmethod
    def get_metadata(file_path: str) -> Dict[str, Any]:
        """Get basic metadata for any file"""
        path = Path(file_path)
        return {
            'filename': path.name,
            'extension': path.suffix.lower(),
            'size_bytes': path.stat().st_size,
            'modified_time': path.stat().st_mtime,
        }

class TextFileHandler(FileHandler):
    """Handler for text-based files"""
    
    # Extensions that are considered text files
    TEXT_EXTENSIONS = {
        '.txt', '.md', '.rst', '.log', '.json', '.xml', '.yaml', '.yml',
        '.ini', '.conf', '.cfg', '.properties',
        '.html', '.htm', '.css', '.scss', '.less',
        '.js', '.ts', '.jsx', '.tsx', '.mjs',
        '.py', '.java', '.c', '.cpp', '.h', '.hpp', '.cs', '.go', '.rb', '.php',
        '.sh', '.bat', '.ps1', '.sql'
    }
    
    @staticmethod
    def can_handle(file_path: str) -> bool:
        """Check if this handler can process the given file"""
        ext = Path(file_path).suffix.lower()
        return ext in TextFileHandler.TEXT_EXTENSIONS and os.path.getsize(file_path) < (5 * 1024 * 1024)  # 5MB limit
    
    @staticmethod
    def extract_text(file_path: str) -> Tuple[str, Dict[str, Any]]:
        """Extract text from a text file"""
        metadata = FileHandler.get_metadata(file_path)
        
        try:
            # Try to detect encoding
            encodings = ['utf-8', 'latin-1', 'cp1252']
            content = ""
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        content = f.read()
                    metadata['encoding'] = encoding
                    break
                except UnicodeDecodeError:
                    continue
            
            if not content:
                logger.warning(f"Could not decode text file: {file_path}")
                return "", metadata
                
            return content, metadata
        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {e}")
            return "", metadata

class PDFHandler(FileHandler):
    """Handler for PDF files"""
    
    @staticmethod
    def can_handle(file_path: str) -> bool:
        """Check if this handler can process the given file"""
        if not PYPDF_AVAILABLE:
            return False
        ext = Path(file_path).suffix.lower()
        return ext == '.pdf' and os.path.getsize(file_path) < (20 * 1024 * 1024)  # 20MB limit
    
    @staticmethod
    def extract_text(file_path: str) -> Tuple[str, Dict[str, Any]]:
        """Extract text from a PDF file"""
        metadata = FileHandler.get_metadata(file_path)
        
        try:
            with open(file_path, 'rb') as f:
                pdf = pypdf.PdfReader(f)
                metadata['page_count'] = len(pdf.pages)
                
                # Extract document info
                if pdf.metadata:
                    for key, value in pdf.metadata.items():
                        if key and value and isinstance(key, str):
                            clean_key = key.strip('/').lower()
                            if isinstance(value, str):
                                metadata[f'pdf_{clean_key}'] = value
                
                # Extract text
                text = ""
                for page in pdf.pages:
                    text += page.extract_text() + "\n\n"
                
                return text, metadata
        except Exception as e:
            logger.error(f"Error extracting text from PDF {file_path}: {e}")
            return "", metadata

class DocxHandler(FileHandler):
    """Handler for DOCX files"""
    
    @staticmethod
    def can_handle(file_path: str) -> bool:
        """Check if this handler can process the given file"""
        if not DOCX_AVAILABLE:
            return False
        ext = Path(file_path).suffix.lower()
        return ext == '.docx' and os.path.getsize(file_path) < (10 * 1024 * 1024)  # 10MB limit
    
    @staticmethod
    def extract_text(file_path: str) -> Tuple[str, Dict[str, Any]]:
        """Extract text from a DOCX file"""
        metadata = FileHandler.get_metadata(file_path)
        
        try:
            doc = Document(file_path)
            
            # Extract document properties
            core_props = doc.core_properties
            if core_props:
                metadata['author'] = core_props.author
                metadata['created'] = core_props.created.isoformat() if core_props.created else None
                metadata['title'] = core_props.title
                
            # Extract text
            text = ""
            for para in doc.paragraphs:
                text += para.text + "\n"
            
            # Get text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text += cell.text + " "
                    text += "\n"
            
            return text, metadata
        except Exception as e:
            logger.error(f"Error extracting text from DOCX {file_path}: {e}")
            return "", metadata

class ExcelHandler(FileHandler):
    """Handler for Excel files"""
    
    @staticmethod
    def can_handle(file_path: str) -> bool:
        """Check if this handler can process the given file"""
        if not EXCEL_AVAILABLE:
            return False
        ext = Path(file_path).suffix.lower()
        return ext in ['.xlsx', '.xls'] and os.path.getsize(file_path) < (10 * 1024 * 1024)  # 10MB limit
    
    @staticmethod
    def extract_text(file_path: str) -> Tuple[str, Dict[str, Any]]:
        """Extract text from an Excel file"""
        metadata = FileHandler.get_metadata(file_path)
        
        try:
            workbook = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            metadata['sheet_names'] = workbook.sheetnames
            
            text = ""
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text += f"Sheet: {sheet_name}\n"
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = []
                    for cell in row:
                        if cell is not None:
                            row_text.append(str(cell))
                    if row_text:
                        text += " | ".join(row_text) + "\n"
                text += "\n"
            
            return text, metadata
        except Exception as e:
            logger.error(f"Error extracting text from Excel {file_path}: {e}")
            return "", metadata

class PowerPointHandler(FileHandler):
    """Handler for PowerPoint files"""
    
    @staticmethod
    def can_handle(file_path: str) -> bool:
        """Check if this handler can process the given file"""
        if not PPTX_AVAILABLE:
            return False
        ext = Path(file_path).suffix.lower()
        return ext in ['.pptx', '.ppt'] and os.path.getsize(file_path) < (15 * 1024 * 1024)  # 15MB limit
    
    @staticmethod
    def extract_text(file_path: str) -> Tuple[str, Dict[str, Any]]:
        """Extract text from a PowerPoint file"""
        metadata = FileHandler.get_metadata(file_path)
        
        try:
            prs = Presentation(file_path)
            metadata['slide_count'] = len(prs.slides)
            
            text = ""
            for i, slide in enumerate(prs.slides):
                text += f"Slide {i+1}:\n"
                
                if slide.title:
                    title_text = slide.title.text if hasattr(slide.title, 'text') else ""
                    text += f"Title: {title_text}\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"
                
                text += "\n"
            
            return text, metadata
        except Exception as e:
            logger.error(f"Error extracting text from PowerPoint {file_path}: {e}")
            return "", metadata

def get_file_handlers() -> List[type]:
    """Get all available file handlers"""
    return [
        TextFileHandler,
        PDFHandler, 
        DocxHandler,
        ExcelHandler,
        PowerPointHandler
    ]

def extract_file_content(file_path: str) -> Tuple[str, Dict[str, Any]]:
    """
    Extract text content from a file using the appropriate handler
    
    Args:
        file_path: Path to the file
        
    Returns:
        Tuple of (extracted text, metadata)
    """
    if not os.path.exists(file_path) or not os.path.isfile(file_path):
        logger.warning(f"File does not exist or is not a regular file: {file_path}")
        return "", {"error": "File not found or not a regular file"}
    
    # Find the appropriate handler
    for handler_class in get_file_handlers():
        if handler_class.can_handle(file_path):
            logger.debug(f"Using {handler_class.__name__} for {file_path}")
            return handler_class.extract_text(file_path)
    
    logger.warning(f"No handler found for file: {file_path}")
    return "", FileHandler.get_metadata(file_path)
